"""
AI-Powered Subsystem Test Plan Generation
Professional presentation — Teal Innovation theme
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
DARK     = RGBColor(0x0A, 0x3D, 0x47)
TEAL     = RGBColor(0x02, 0x80, 0x90)
SEAFOAM  = RGBColor(0x00, 0xA8, 0x96)
MINT     = RGBColor(0x02, 0xC3, 0x9A)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHTBG  = RGBColor(0xF0, 0xF9, 0xFA)
DARKTEXT = RGBColor(0x1A, 0x2A, 0x3A)
MUTED    = RGBColor(0x5A, 0x7A, 0x85)
ACCENT   = RGBColor(0xFF, 0x6B, 0x35)
CODEBG   = RGBColor(0x0D, 0x11, 0x17)
CODEGRAY = RGBColor(0x8B, 0x94, 0x9E)
CODEBLUE = RGBColor(0x79, 0xC0, 0xFF)
CODELTBL = RGBColor(0xA5, 0xD6, 0xFF)
CODEFG   = RGBColor(0xE6, 0xED, 0xF3)
RED      = RGBColor(0xC0, 0x39, 0x2B)
REDBG    = RGBColor(0xFF, 0xF5, 0xF5)
GREEN    = RGBColor(0x27, 0xAE, 0x60)
GREENBG  = RGBColor(0xF0, 0xFF, 0xF4)
ORANGE   = RGBColor(0xE6, 0x7E, 0x22)
ORANGEBG = RGBColor(0xFF, 0xFA, 0xF0)
PURPLE   = RGBColor(0x8E, 0x44, 0xAD)
PURPLEBG = RGBColor(0xFA, 0xF0, 0xFF)
TEALBG   = RGBColor(0xF0, 0xFA, 0xFA)
SEAFOAMBG= RGBColor(0xF0, 0xFE, 0xFA)
CARDWHT  = RGBColor(0xFF, 0xFF, 0xFF)
CARDBRD  = RGBColor(0xE0, 0xE0, 0xE0)
STRIPALT = RGBColor(0xF5, 0xFE, 0xFF)
HEADERLT = RGBColor(0xF7, 0xFE, 0xFF)

W = 10.0   # slide width  inches
H = 5.625  # slide height inches

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)

# ── Low-level shape helpers ──────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width_pt=0):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt) if line_width_pt else Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_line_shape(slide, x1, y1, length, fill_color):
    """Horizontal line using a very thin rectangle."""
    add_rect(slide, x1, y1, length, 0.02, fill_color)

def add_text(slide, text, x, y, w, h,
             size=11, bold=False, italic=False, color=None,
             align=PP_ALIGN.LEFT, valign=None, font="Calibri",
             wrap=True, bg=None, line_color=None):
    from pptx.enum.text import MSO_ANCHOR
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if valign:
        tf.auto_size = None
        tf.vertical_anchor = valign
    if bg:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg
    else:
        txBox.fill.background()
    txBox.line.fill.background()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    if color:
        run.font.color.rgb = color
    return txBox

def add_rich_text(slide, runs, x, y, w, h,
                  size=11, align=PP_ALIGN.LEFT, valign=None,
                  wrap=True, bg=None, line_color=None, para_space_after=None):
    """runs = list of (text, bold, italic, color, size_override, font_override)"""
    from pptx.enum.text import MSO_ANCHOR
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if valign:
        tf.auto_size = None
        tf.vertical_anchor = valign
    if bg:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg
    else:
        txBox.fill.background()
    txBox.line.fill.background()
    # First paragraph already exists
    first = True
    cur_para = tf.paragraphs[0]
    cur_para.alignment = align
    for item in runs:
        text = item[0]
        bold    = item[1] if len(item) > 1 else False
        italic  = item[2] if len(item) > 2 else False
        color   = item[3] if len(item) > 3 else None
        sz      = item[4] if len(item) > 4 else size
        font    = item[5] if len(item) > 5 else "Calibri"
        if text == "\n":
            cur_para = tf.add_paragraph()
            cur_para.alignment = align
            if para_space_after:
                from pptx.util import Pt as _Pt
                cur_para.space_after = _Pt(para_space_after)
            continue
        run = cur_para.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        if color:
            run.font.color.rgb = color
    return txBox

def add_bullet_list(slide, items, x, y, w, h, size=10, color=None,
                    indent_color=None, bold_first=False, font="Calibri"):
    """Items: list of strings. Each prefixed with bullet dot."""
    from pptx.oxml.ns import qn
    from lxml import etree
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    txBox.fill.background()
    txBox.line.fill.background()
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "• " + item
        run.font.size = Pt(size)
        run.font.name = font
        run.font.bold = bold_first and (i == 0)
        if color:
            run.font.color.rgb = color
    return txBox

# ── Slide-level helpers ──────────────────────────────────────────────────────
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header_bar(slide, title, subtitle=""):
    """Dark top bar + mint underline + title text."""
    add_rect(slide, 0, 0, W, 0.68, DARK)
    add_rect(slide, 0, 0.68, W, 0.05, MINT)
    add_text(slide, title, 0.35, 0.05, W-0.5, 0.58,
             size=20, bold=True, color=WHITE, font="Calibri")
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.72, W-0.5, 0.32,
                 size=10, italic=True, color=MUTED, font="Calibri")

def add_card(slide, x, y, w, h, bg=None, border=None, accent_color=None):
    bg_c  = bg     or CARDWHT
    brd_c = border or CARDBRD
    add_rect(slide, x, y, w, h, bg_c, brd_c, 0.75)
    if accent_color:
        add_rect(slide, x, y, 0.055, h, accent_color)

def add_section_header(slide, title, x, y, w, h, bg_color):
    add_rect(slide, x, y, w, h, bg_color)
    add_text(slide, title, x, y, w, h,
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
def slide_title(prs):
    layout = prs.slide_layouts[6]  # blank
    sl = prs.slides.add_slide(layout)
    set_bg(sl, DARK)

    # Right decorative columns
    add_rect(sl, 7.6, 0, 2.4, H, TEAL)
    add_rect(sl, 8.5, 0, 1.5, H, RGBColor(0x01, 0x6E, 0x7D))
    add_rect(sl, 9.2, 0, 0.8, H, RGBColor(0x01, 0x5F, 0x6C))

    # Left mint accent
    add_rect(sl, 0, 0, 0.09, H, MINT)

    # Main title
    add_rich_text(sl,
        [
            ("AI-Powered Subsystem\n", True, False, WHITE, 34, "Calibri"),
            ("Test Plan Generation", True, False, WHITE, 34, "Calibri"),
        ],
        0.45, 1.0, 6.8, 2.0)

    # Divider
    add_rect(sl, 0.45, 3.1, 6.5, 0.03, SEAFOAM)

    # Subtitle
    add_text(sl,
        "Using ChipAgent  ·  From Design Specification to Structured Test Plan",
        0.45, 3.22, 6.8, 0.45, size=13, italic=True, color=MINT)

    # Tags
    tags = ["ChipAgent", "Test Automation", "AI Skill", "Design Spec"]
    for i, tag in enumerate(tags):
        tx = 0.45 + i * 1.58
        add_rect(sl, tx, 3.9, 1.45, 0.38, TEAL, SEAFOAM, 0.75)
        add_text(sl, tag, tx, 3.9, 1.45, 0.38,
                 size=10, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Planning the Skill
# ════════════════════════════════════════════════════════════════════════════
def slide_planning(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, LIGHTBG)
    add_header_bar(sl, "Planning the Skill",
                   "Init Prompt  ·  Key Questions  ·  Use Skill Creator")

    # ── Left: Init prompt block ──
    px, py, pw, ph = 0.28, 1.1, 5.55, 4.22
    add_card(sl, px, py, pw, ph, accent_color=TEAL)
    add_text(sl, "Sample Init Prompt", px+0.14, py+0.08, pw-0.2, 0.35,
             size=12, bold=True, color=DARK)

    # Prompt code-style box
    add_rect(sl, px+0.1, py+0.5, pw-0.18, ph-0.62, HEADERLT, CARDBRD, 0.5)
    prompt_runs = [
        ('"I want to build a ', False, False, DARKTEXT, 10.5),
        ("ChipAgent skill", True, False, TEAL, 10.5),
        (" to automatically generate a ", False, False, DARKTEXT, 10.5),
        ("subsystem test plan", True, False, TEAL, 10.5),
        (" from a ", False, False, DARKTEXT, 10.5),
        ("design specification document", True, False, TEAL, 10.5),
        (".", False, False, DARKTEXT, 10.5),
        ("\n", ),
        ("What I know so far:", True, False, DARKTEXT, 10.5),
        ("\n",),
        ("• Input: Design spec (PDF / DOCX / structured text)", False, False, DARKTEXT, 10),
        ("\n",),
        ("• Output: Structured test plan in Excel / Word format", False, False, DARKTEXT, 10),
        ("\n",),
        ("• Skill must be modular, reusable, token-optimised", False, False, DARKTEXT, 10),
        ("\n",),
        ("• Benchmark against existing plans, MDA agent, prior versions", False, False, DARKTEXT, 10),
        ("\n",),
        ("Please ask me questions to fill any gaps, then use ", False, False, DARKTEXT, 10.5),
        ("Skill Creator", True, False, ACCENT, 10.5),
        (' to create this skill."', False, False, DARKTEXT, 10.5),
    ]
    add_rich_text(sl, prompt_runs, px+0.18, py+0.55, pw-0.32, ph-0.72,
                  size=10.5, wrap=True)

    # ── Right: two info cards ──
    rx = 6.0

    # Card A — AI will ask about
    add_card(sl, rx, 1.1, 3.72, 1.92, accent_color=SEAFOAM)
    add_text(sl, "AI Will Ask About", rx+0.14, 1.15, 3.4, 0.35,
             size=12, bold=True, color=DARK)
    add_bullet_list(sl, [
        "Document format (PDF, DOCX, plain text)",
        "Output format / Excel template columns",
        "Subsystem types & scope boundaries",
        "Benchmarking material availability",
    ], rx+0.14, 1.55, 3.45, 1.35, size=10, color=DARKTEXT)

    # Card B — also include
    add_card(sl, rx, 3.15, 3.72, 2.12, accent_color=MINT)
    add_text(sl, "Also Include in Your Init Prompt", rx+0.14, 3.2, 3.45, 0.35,
             size=12, bold=True, color=DARK)
    add_bullet_list(sl, [
        "Existing test plan samples (reference)",
        "Expected sections: objectives, test cases, coverage",
        "Acceptance criteria / quality bar",
        "Known edge cases or special signal types",
        "Team conventions: naming, IDs, priority levels",
    ], rx+0.14, 3.6, 3.45, 1.55, size=10, color=DARKTEXT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Inputs · Outputs · Benchmarking
# ════════════════════════════════════════════════════════════════════════════
def slide_io(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, LIGHTBG)
    add_header_bar(sl, "Skill Definition: Inputs · Outputs · Benchmarking")

    cols = [
        {
            "title": "Inputs to the Skill",
            "color": TEAL,
            "x": 0.22,
            "items": [
                "Design specification document\n(PDF / DOCX / structured text)",
                "Subsystem name & scope",
                "Output Excel/Word template",
                "Reference test plan (existing)",
                "Protocol / standard reference docs",
                "Parsing config: section headers,\nnaming conventions",
            ]
        },
        {
            "title": "Output & Format",
            "color": SEAFOAM,
            "x": 3.52,
            "items": [
                "Structured test plan document",
                "Excel: Test ID, description,\nexpected result, priority, tag",
                "Word: organised by feature /\nsubsystem section",
                "Intermediate JSON/YAML\n(for downstream skills)",
                "Interface & protocol summary\n(reusable artefact)",
                "Parsing log / section map",
            ]
        },
        {
            "title": "Benchmarking",
            "color": ACCENT,
            "x": 6.82,
            "items": [
                "vs. Existing test plan\n(coverage %, section match)",
                "vs. MDA Agent output\n(quality & completeness)",
                "vs. Previous skill version\n(regression & improvement)",
                "Metrics: # test cases, coverage\ndepth, missed requirements",
                "False-positive test cases flagged",
                "Manual review checklist",
            ]
        }
    ]

    col_w = 3.08
    row_h = 0.65
    for col in cols:
        x = col["x"]
        # Section header
        add_rect(sl, x, 1.0, col_w, 0.42, col["color"])
        add_text(sl, col["title"], x, 1.0, col_w, 0.42,
                 size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font="Calibri")
        # Rows
        for i, item in enumerate(col["items"]):
            ry = 1.52 + i * row_h
            bg = CARDWHT if i % 2 == 0 else STRIPALT
            add_rect(sl, x, ry, col_w, row_h - 0.04, bg, CARDBRD, 0.5)
            add_rect(sl, x, ry, 0.05, row_h - 0.04, col["color"])
            add_text(sl, item, x+0.1, ry+0.04, col_w-0.14, row_h-0.1,
                     size=9.5, color=DARKTEXT, wrap=True, font="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Skill Structure
# ════════════════════════════════════════════════════════════════════════════
def slide_structure(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, LIGHTBG)
    add_header_bar(sl, "Structure of the Skill",
                   "Modular · Reusable · Generic · Token-Optimised")

    # ── Principles row ──
    principles = [
        ("Modular",         "Each step is an\nindependent sub-skill",    TEAL),
        ("Reusable",        "Scripts shared across\nmultiple skills",      SEAFOAM),
        ("Generic",         "Works across different\nsubsystem types",     MINT),
        ("Token-Optimised", "Compact context,\nno redundant steps",        RGBColor(0x02,0x70,0x80)),
    ]
    for i, (lbl, desc, col) in enumerate(principles):
        px = 0.28 + i * 2.38
        add_rect(sl, px, 1.05, 2.22, 0.75, col)
        add_rich_text(sl,
            [(lbl+"\n", True, False, WHITE, 12, "Calibri"),
             (desc,     False, True,  WHITE, 9.5, "Calibri")],
            px, 1.05, 2.22, 0.75,
            align=PP_ALIGN.CENTER)

    # ── Sub-skills table ──
    add_card(sl, 0.28, 1.92, 5.9, 3.42, accent_color=TEAL)
    add_text(sl, "Sub-Skills & Intermediate Outputs", 0.42, 1.96, 5.6, 0.34,
             size=12, bold=True, color=DARK)

    sub_skills = [
        ("doc-parser",          "Parses design spec → section map JSON",           "Reused by interface-extractor & VIP dev skill"),
        ("req-extractor",       "Extracts requirements & features from sections",   "Feeds test-case-generator"),
        ("interface-extractor", "Pulls interface & protocol info",                  "Used by Interface VIP Dev Skill"),
        ("testcase-generator",  "Generates test cases per requirement",             "Core output of this skill"),
        ("xls-writer",          "Writes test plan to Excel template",               "xls-reader script reused by other skills"),
    ]
    for i, (name, desc, downstream) in enumerate(sub_skills):
        sy = 2.38 + i * 0.56
        # Name pill
        add_rect(sl, 0.35, sy, 1.55, 0.42, TEAL)
        add_text(sl, name, 0.35, sy, 1.55, 0.42,
                 size=8.5, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font="Consolas")
        # Arrow
        add_rect(sl, 1.93, sy+0.19, 0.18, 0.02, SEAFOAM)
        # Description
        add_text(sl, desc, 2.14, sy+0.03, 2.3, 0.36, size=9.5, color=DARKTEXT)
        # Downstream
        add_text(sl, "→ "+downstream, 4.47, sy+0.03, 1.6, 0.36,
                 size=8.5, italic=True, color=TEAL)

    # ── What else ──
    add_card(sl, 6.32, 1.92, 3.45, 3.42, accent_color=SEAFOAM)
    add_text(sl, "What Else Belongs Here", 6.46, 1.96, 3.15, 0.34,
             size=12, bold=True, color=DARK)
    add_bullet_list(sl, [
        "Error-recovery scripts\n  (retry logic, partial-result save)",
        "Context-reset utilities\n  (/clear + /compact mid-task)",
        "Validation / QA scripts\n  (coverage checker, diff vs reference)",
        "Shared Excel template\n  (standard test-plan format)",
        "Glossary / abbreviation resolver\n  (domain terms for LLM context)",
    ], 6.46, 2.38, 3.18, 2.82, size=9.5, color=DARKTEXT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Good Practices
# ════════════════════════════════════════════════════════════════════════════
def slide_practices(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, LIGHTBG)
    add_header_bar(sl, "Good Practices",
                   "Quality · Monitoring · Environment · Infrastructure")

    cards = [
        ("Stop on Error — Don't Improvise",
         "Instruct agent to halt & report immediately on any error.\n"
         "Do NOT find workarounds — they silently compromise output quality.",
         RED, REDBG),
        ("Monitor Processing Closely",
         "Watch for: errors, context approaching 100%, slow/looping steps.\n"
         "Embed repetitive non-core steps (e.g. doc parsing) as scripts\nto preserve token budget.",
         TEAL, TEALBG),
        ("Use Local Virtual Environment",
         "Always use a Python venv for all dependencies.\n"
         "Never install globally — ensures reproducibility & isolation.",
         GREEN, GREENBG),
        ("Keep Workspace Clean",
         "Don't clutter workspace with intermediate docs not needed for the task.\n"
         "Remove temp files after each step to reduce noise & context usage.",
         ORANGE, ORANGEBG),
        ("Invest in Common Infrastructure",
         "Shared scripts (xls-reader, doc-parser), common templates & standard\n"
         "spec formats reduce AI pressure, cut tokens, improve consistency.",
         SEAFOAM, SEAFOAMBG),
        ("Use New / Clear / Compact Between Tasks",
         "Reset context between major task phases (/new, /clear, /compact).\n"
         "Use as transition points to keep context clean, fresh, and focused.",
         PURPLE, PURPLEBG),
    ]

    cw, ch = 3.08, 2.08
    for i, (title, body, col, bg) in enumerate(cards):
        col_i = i % 3
        row_i = i // 3
        cx = 0.22 + col_i * 3.22
        cy = 1.1  + row_i * 2.18

        add_rect(sl, cx, cy, cw, ch, bg, CARDBRD, 0.5)
        add_rect(sl, cx, cy, cw, 0.055, col)  # top accent bar

        add_text(sl, title, cx+0.12, cy+0.1, cw-0.2, 0.4,
                 size=11, bold=True, color=col, font="Calibri")
        add_text(sl, body, cx+0.12, cy+0.55, cw-0.2, 1.4,
                 size=9.5, color=DARKTEXT, wrap=True, font="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ChipAgent Configuration
# ════════════════════════════════════════════════════════════════════════════
def slide_config(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, LIGHTBG)
    add_header_bar(sl, "ChipAgent Configuration",
                   "Skills · Automation · Rules · Session Triggers")

    cfg_blocks = [
        {
            "title": "Enable Skills",
            "color": TEAL,
            "x": 0.22,
            "lines": [
                ("# Register in skills registry",    CODEGRAY),
                ("skills:",                           CODEBLUE),
                ("  - name: subsystem-testplan",      CODELTBL),
                ("    path: ./skills/testplan",       CODELTBL),
                ("    version: 1.0",                  CODELTBL),
                ("enable-skills: true",               CODEFG),
                ("auto-load: true",                   CODEFG),
            ]
        },
        {
            "title": "Automatic Commands",
            "color": SEAFOAM,
            "x": 3.55,
            "lines": [
                ("# Auto-run on session start",       CODEGRAY),
                ("on_start:",                         CODEBLUE),
                ("  - /load-skill testplan",          CODELTBL),
                ("  - /set-context subsystem",        CODELTBL),
                ("# On trigger keyword",              CODEGRAY),
                ('on_trigger: "generate testplan"',   CODEBLUE),
                ("  - /run testplan-skill",           CODELTBL),
            ]
        },
        {
            "title": "Rules (Always Applied)",
            "color": ACCENT,
            "x": 6.88,
            "lines": [
                ("# Applied every session / trigger", CODEGRAY),
                ("rules:",                            CODEBLUE),
                ("  - stop_on_error: true",           CODELTBL),
                ("  - use_venv: ./venv",              CODELTBL),
                ("  - clear_workspace_after_task",    CODELTBL),
                ("  - context_limit: compact",        CODELTBL),
                ("  - output_format: excel+json",     CODELTBL),
            ]
        }
    ]

    cw = 3.1
    for cfg in cfg_blocks:
        x = cfg["x"]
        # Header
        add_rect(sl, x, 1.0, cw, 0.42, cfg["color"])
        add_text(sl, cfg["title"], x, 1.0, cw, 0.42,
                 size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font="Calibri")
        # Code background
        add_rect(sl, x, 1.46, cw, 3.72, CODEBG, RGBColor(0x30,0x36,0x3D), 0.75)
        # Code lines
        for i, (line, col) in enumerate(cfg["lines"]):
            add_text(sl, line, x+0.12, 1.56+i*0.48, cw-0.18, 0.42,
                     size=9.5, color=col, font="Consolas")

    # Footer note
    add_text(sl,
        "💡  Rules defined in ChipAgent config persist across all sessions — no need to repeat instructions each time.",
        0.22, 5.2, 9.55, 0.32,
        size=9.5, italic=True, color=MUTED)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Iteration to Improve the Skill
# ════════════════════════════════════════════════════════════════════════════
def slide_iteration(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, LIGHTBG)
    add_header_bar(sl, "Iteration to Improve the Skill",
                   "Benchmark → Review → Feedback → Repeat")

    # ── Cycle steps ──
    steps = [
        ("1", "Run Skill &\nBenchmark",   TEAL),
        ("2", "Compare vs\nExisting Plan", SEAFOAM),
        ("3", "Identify\nGaps & Issues",  ACCENT),
        ("4", "Give Feedback\nto Skill",  PURPLE),
        ("5", "Update Skill\n& Re-run",   GREEN),
    ]
    box_w = 1.7
    for i, (num, lbl, col) in enumerate(steps):
        bx = 0.3 + i * 1.88
        # Circle
        circle = sl.shapes.add_shape(9, Inches(bx+0.4), Inches(1.0),
                                     Inches(0.9), Inches(0.9))
        circle.fill.solid()
        circle.fill.fore_color.rgb = col
        circle.line.fill.background()
        add_text(sl, num, bx+0.4, 1.0, 0.9, 0.9,
                 size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font="Calibri")
        # Box
        add_rect(sl, bx, 2.05, box_w, 0.72, CARDWHT, col, 1.5)
        add_text(sl, lbl, bx, 2.05, box_w, 0.72,
                 size=10, bold=True, color=col,
                 align=PP_ALIGN.CENTER, font="Calibri")
        # Arrow
        if i < len(steps) - 1:
            add_rect(sl, bx+box_w+0.02, 2.38, 0.15, 0.02,
                     RGBColor(0xAA,0xAA,0xAA))

    # ── Repeat banner ──
    add_rect(sl, 0.3, 2.9, 9.4, 0.34, DARK)
    add_text(sl,
        "Repeat this cycle multiple times — each iteration makes the skill measurably better",
        0.3, 2.9, 9.4, 0.34,
        size=11, italic=True, color=MINT,
        align=PP_ALIGN.CENTER, font="Calibri")

    # ── Two detail cards ──
    add_card(sl, 0.3, 3.35, 4.55, 2.0, accent_color=TEAL)
    add_text(sl, "What to Review in Each Iteration",
             0.44, 3.4, 4.3, 0.34, size=11, bold=True, color=DARK)
    add_bullet_list(sl, [
        "Coverage %: requirements with test cases",
        "Missing scenarios flagged by reviewer",
        "Incorrect or vague test descriptions",
        "Section structure match vs existing plan",
        "Delta improvement vs prior skill version",
    ], 0.44, 3.8, 4.28, 1.45, size=9.5, color=DARKTEXT)

    add_card(sl, 5.05, 3.35, 4.65, 2.0, accent_color=SEAFOAM)
    add_text(sl, "How to Give Effective Feedback",
             5.19, 3.4, 4.4, 0.34, size=11, bold=True, color=DARK)
    add_bullet_list(sl, [
        'Be specific: "Interrupt handling test missing"',
        "Annotate generated plan before feeding back",
        "Compare diff vs MDA output for systematic gaps",
        "Update skill prompt after each review cycle",
        "Track skill version history with benchmark scores",
    ], 5.19, 3.8, 4.38, 1.45, size=9.5, color=DARKTEXT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Known Issues & Mitigations
# ════════════════════════════════════════════════════════════════════════════
def slide_issues(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, LIGHTBG)
    add_header_bar(sl, "Known Issues & Mitigations",
                   "Real challenges observed while running ChipAgent skills in production")

    issues = [
        {
            "title": "Context Stuck at Fixed %",
            "icon": "🔒",
            "symptom": (
                "ChipAgent context meter freezes at a fixed percentage (e.g. 72%) "
                "mid-task and stops progressing — the agent appears to keep running "
                "but no meaningful output is produced."
            ),
            "root_cause": (
                "Large intermediate documents or repeated tool calls inflate the "
                "context without advancing the task. The agent enters a loop trying "
                "to fit work into an already-saturated context window."
            ),
            "mitigation": [
                "Monitor context % continuously — set an alert threshold (e.g. 80%)",
                "Use /compact or /clear at natural task boundaries before hitting the limit",
                "Break large documents into chunks; process one section at a time",
                "Store intermediate outputs to disk — don't keep them in context",
                "Embed heavy parsing steps as scripts so they don't consume context",
            ],
            "color": RED,
            "bg": REDBG,
            "x": 0.25,
        },
        {
            "title": "Agent Avoids Reporting Errors",
            "icon": "⚠️",
            "symptom": (
                "When ChipAgent encounters an error (file not found, parse failure, "
                "unexpected format), it silently tries an alternative path instead of "
                "stopping — often producing low-quality or incorrect output."
            ),
            "root_cause": (
                "Default agent behaviour is to be 'helpful' and keep going. Without "
                "an explicit stop-on-error rule, the agent treats errors as obstacles "
                "to route around rather than signals to halt and report."
            ),
            "mitigation": [
                "Add stop_on_error: true to ChipAgent rules config",
                "Include in skill prompt: 'If you hit any error, stop immediately and report it'",
                "Add a validation step after each sub-skill — check outputs before proceeding",
                "Review agent logs after every run — errors often hide in verbose output",
                "Use structured output (JSON/YAML) so missing fields are detectable",
            ],
            "color": ORANGE,
            "bg": ORANGEBG,
            "x": 5.1,
        },
    ]

    for iss in issues:
        x, w, h = iss["x"], 4.6, 4.15
        y = 1.08
        add_rect(sl, x, y, w, h, iss["bg"], CARDBRD, 0.75)
        add_rect(sl, x, y, w, 0.055, iss["color"])

        # Issue title
        add_text(sl, iss["icon"] + "  " + iss["title"],
                 x+0.14, y+0.1, w-0.2, 0.38,
                 size=13, bold=True, color=iss["color"], font="Calibri")

        # Symptom
        add_rect(sl, x+0.12, y+0.55, w-0.22, 0.78,
                 RGBColor(0xFF,0xFF,0xFF), CARDBRD, 0.5)
        add_text(sl, "Symptom:", x+0.2, y+0.57, w-0.3, 0.2,
                 size=9.5, bold=True, color=iss["color"], font="Calibri")
        add_text(sl, iss["symptom"], x+0.2, y+0.78, w-0.3, 0.5,
                 size=9, color=DARKTEXT, wrap=True, font="Calibri")

        # Root cause
        add_rect(sl, x+0.12, y+1.38, w-0.22, 0.68,
                 RGBColor(0xFF,0xFF,0xFF), CARDBRD, 0.5)
        add_text(sl, "Root Cause:", x+0.2, y+1.4, w-0.3, 0.2,
                 size=9.5, bold=True, color=DARKTEXT, font="Calibri")
        add_text(sl, iss["root_cause"], x+0.2, y+1.6, w-0.3, 0.42,
                 size=9, color=DARKTEXT, wrap=True, font="Calibri")

        # Mitigations
        add_text(sl, "Mitigations:", x+0.14, y+2.12, w-0.2, 0.26,
                 size=10, bold=True, color=iss["color"], font="Calibri")
        add_bullet_list(sl, iss["mitigation"],
                        x+0.14, y+2.38, w-0.2, 1.65,
                        size=9.5, color=DARKTEXT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Key Takeaways (dark close)
# ════════════════════════════════════════════════════════════════════════════
def slide_takeaways(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, DARK)

    # Left mint bar
    add_rect(sl, 0, 0, 0.1, H, MINT)
    # Bottom teal bar
    add_rect(sl, 0, H-0.14, W, 0.14, TEAL)

    add_text(sl, "Key Takeaways", 0.45, 0.3, 9.0, 0.65,
             size=26, bold=True, color=WHITE, font="Calibri")
    add_rect(sl, 0.45, 1.0, 8.6, 0.03, MINT)

    takeaways = [
        ("01", "A well-planned init prompt + Skill Creator = faster, higher-quality skill development",   TEAL),
        ("02", "Modular sub-skills create reusable building blocks across the entire ChipAgent ecosystem", SEAFOAM),
        ("03", "Good practices (stop on error, clean workspace, venv, compact) ensure consistent quality", MINT),
        ("04", "Multiple benchmark → review → feedback cycles are essential for production-quality skills",ACCENT),
    ]
    for i, (num, txt, col) in enumerate(takeaways):
        ty = 1.12 + i * 0.98
        add_rect(sl, 0.45, ty, 0.58, 0.58, col)
        add_text(sl, num, 0.45, ty, 0.58, 0.58,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font="Calibri")
        add_text(sl, txt, 1.18, ty+0.08, 8.5, 0.48,
                 size=13, color=WHITE, font="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# BUILD
# ════════════════════════════════════════════════════════════════════════════
slide_title(prs)
slide_planning(prs)
slide_io(prs)
slide_structure(prs)
slide_practices(prs)
slide_config(prs)
slide_iteration(prs)
slide_issues(prs)
slide_takeaways(prs)

OUT = "/Users/apple/Documents/project-ppt/AI_TestPlan_Skill.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
